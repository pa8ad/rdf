#!/usr/bin/env python3
import os
import pandas as pd
from pathlib import Path
from datetime import datetime
import hashlib
import magic
from PyPDF2 import PdfReader
from docx import Document
import openpyxl
from pptx import Presentation
import warnings
from dotenv import load_dotenv
from openai import OpenAI
import base64
import json
import tkinter as tk
from tkinter import filedialog, messagebox

# Load environment variables from custom .env file
load_dotenv('sleu.env')

# Suppress openpyxl warnings
warnings.filterwarnings(
    "ignore",
    message="Print area cannot be set to Defined name:.*",
    category=UserWarning,
    module="openpyxl.reader.workbook"
)

# Initialize OpenAI client
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise ValueError("Geen OpenAI API-sleutel gevonden. Zet OPENAI_API_KEY in sleu.env.")
client = OpenAI(api_key=api_key)


def compute_checksum(filepath, algorithm='sha256', chunk_size=8192):
    h = hashlib.new(algorithm)
    with open(filepath, 'rb') as f:
        for chunk in iter(lambda: f.read(chunk_size), b''):
            h.update(chunk)
    return h.hexdigest()


def generate_summary(text, word_limit=100):
    snippet = text[:3000]
    messages = [
        {"role": "system", "content": "Je bent een behulpzame AI-assistent gespecialiseerd in document-samenvattingen."},
        {"role": "user", "content": f"Vat het volgende document samen in ongeveer {word_limit} woorden:\n\n{snippet}"}
    ]
    try:
        resp = client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            max_tokens=300,
            temperature=0.5,
        )
        return resp.choices[0].message.content.strip()
    except Exception as e:
        return f"[Samenvatting kon niet worden gegenereerd: {e}]"


def generate_image_description(filepath, word_limit=100):
    """
    Generate a description for an image file using a vision-enabled model,
    limited to approximately word_limit words.
    """
    try:
        with open(filepath, 'rb') as f:
            img = f.read()
        mime_type = 'image/jpeg'
        b64 = base64.b64encode(img).decode()
        data_url = f"data:{mime_type};base64,{b64}"
        prompt = (
            f"Beschrijf in het Nederlands wat je op deze afbeelding ziet in ongeveer {word_limit} woorden."
        )
        messages = [
            {"role": "system", "content": "Je bent een behulpzame AI-assistent gespecialiseerd in beeldbeschrijvingen."},
            {"role": "user", "content": prompt},
            {"role": "user", "content": [
                {"type": "image_url", "image_url": {"url": data_url}}
            ]}
        ]
        resp = client.chat.completions.create(
            model="gpt-4o",
            messages=messages
        )
        desc = resp.choices[0].message.content.strip()
        words = desc.split()
        if len(words) > word_limit:
            desc = ' '.join(words[:word_limit])
        return desc
    except Exception as e:
        return f"[Beschrijving kon niet worden gegenereerd: {e}]"


def extract_metadata(directory, include_subdirs=True):
    records = []
    mime_detector = magic.Magic(mime=True)

    for root, dirs, files in os.walk(directory):
        if not include_subdirs:
            dirs[:] = []
        for fname in files:
            path = Path(root) / fname
            try:
                stat = path.stat()
            except OSError:
                continue

            meta = {
                'file_path': str(path.parent.resolve()),
                'size_bytes': stat.st_size,
                'created_timestamp': datetime.fromtimestamp(stat.st_ctime).strftime('%Y-%m-%d %H:%M:%S'),
                'modified_timestamp': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
                'permissions_mode': stat.st_mode,
                'owner_uid': stat.st_uid,
                'group_gid': stat.st_gid,
                'mime_type': mime_detector.from_file(str(path)),
                'file_extension': path.suffix.lower(),
                'md5_checksum': compute_checksum(str(path), 'md5'),
                'sha256_checksum': compute_checksum(str(path), 'sha256')
            }
            text = ''
            ext = path.suffix.lower()

            if ext == '.pdf':
                reader = PdfReader(str(path))
                info = reader.metadata
                meta.update({
                    'pdf_num_pages': len(reader.pages),
                    'pdf_author': info.author,
                    'pdf_title': info.title
                })
                for pg in reader.pages:
                    text += pg.extract_text() or ''
                meta['pdf_word_count'] = len(text.split())

            elif ext == '.docx':
                doc = Document(str(path))
                props = doc.core_properties
                paras = [p.text for p in doc.paragraphs if p.text]
                text = '\n'.join(paras)
                meta.update({
                    'docx_author': props.author,
                    'docx_title': props.title,
                    'docx_created': props.created.strftime('%Y-%m-%d %H:%M:%S') if props.created else None,
                    'docx_modified': props.modified.strftime('%Y-%m-%d %H:%M:%S') if props.modified else None,
                    'docx_paragraph_count': len(paras),
                    'docx_word_count': len(text.split())
                })

            elif ext == '.pptx':
                prs = Presentation(str(path))
                props = prs.core_properties
                slides = len(prs.slides)
                boxes = 0
                for sl in prs.slides:
                    for shp in sl.shapes:
                        if hasattr(shp, 'text') and shp.text:
                            text += shp.text + '\n'
                            boxes += 1
                meta.update({
                    'pptx_author': props.author,
                    'pptx_title': props.title,
                    'pptx_created': props.created.strftime('%Y-%m-%d %H:%M:%S') if props.created else None,
                    'pptx_modified': props.modified.strftime('%Y-%m-%d %H:%M:%S') if props.modified else None,
                    'pptx_slide_count': slides,
                    'pptx_textbox_count': boxes
                })

            elif ext in ['.xlsx', '.xls']:
                try:
                    wb = openpyxl.load_workbook(str(path), data_only=True)
                    sheets = wb.sheetnames
                    cells = []
                    for sh in sheets:
                        ws = wb[sh]
                        for row in ws.iter_rows(values_only=True):
                            for c in row:
                                if isinstance(c, str):
                                    cells.append(c)
                    text = '\n'.join(cells)
                    meta.update({
                        'excel_sheets': sheets,
                        'excel_word_count': len(text.split())
                    })
                except Exception:
                    pass

            elif ext in ['.jpg', '.jpeg']:
                meta['image_description'] = generate_image_description(str(path))

            if text:
                meta['ai_summary'] = generate_summary(text, word_limit=100)

            for k, v in meta.items():
                records.append({'file_key': path.name, 'metadata_field': k, 'metadata_value': v})

    return pd.DataFrame(records)


def main():
    root = tk.Tk()
    root.withdraw()

    directory = filedialog.askdirectory(title="Kies de map met documenten")
    if not directory:
        print("Geen map geselecteerd. Afsluiten.")
        return

    output = filedialog.asksaveasfilename(
        title="Opslaan als",
        defaultextension=".csv",
        filetypes=[("CSV files", "*.csv")],
        initialfile="metadata.csv"
    )
    if not output:
        print("Geen output-bestand geselecteerd. Afsluiten.")
        return

    include_subdirs = messagebox.askyesno("Submappen", "Submappen meenemen?")

    # Show progress dialog
    progress = tk.Toplevel(root)
    progress.title("Bezig met verwerken")
    tk.Label(progress, text="Metadata-extractie gestart...\nEven geduld alstublieft.").pack(padx=20, pady=20)
    progress.update()

    df = extract_metadata(directory, include_subdirs=include_subdirs)
    df.to_csv(output, index=False)

    # Schrijf MDTO-format JSON
    mdto_path = Path(output).with_suffix('.mdto.json')
    grouped = df.groupby('file_key')
    files = []
    for key, grp in grouped:
        meta_dict = dict(zip(grp['metadata_field'], grp['metadata_value']))
        files.append({'file_key': key, 'metadata': meta_dict})
    with open(mdto_path, 'w', encoding='utf-8') as jf:
        json.dump({'files': files}, jf, ensure_ascii=False, indent=2)

    progress.destroy()
    messagebox.showinfo(
        "Klaar",
        f"Alle metadata is opgeslagen in {output} (CSV) en {mdto_path} (MDTO-format)\n"
        f"Submappen {'inclusief' if include_subdirs else 'exclusief'}."
    )
    root.destroy()


if __name__ == '__main__':
    main()
